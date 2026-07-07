import base64
import io
import json
import os
import sys

from pptx import Presentation
from PIL import Image

if len(sys.argv) < 3:
    print("Uso: python gerar_balanco_pptx.py <input.json> <output.pptx>", file=sys.stderr)
    sys.exit(1)

input_json = sys.argv[1]
output_pptx = sys.argv[2]

with open(input_json, encoding="utf-8") as f:
    data = json.load(f)

script_dir = os.path.dirname(os.path.abspath(__file__))
project_root = os.path.dirname(script_dir)
template_path = os.path.join(project_root, "public", "templates", "Modelo_de_PPT.pptx")

if not os.path.isfile(template_path):
    print(f"Template não encontrado: {template_path}", file=sys.stderr)
    sys.exit(1)

prs = Presentation(template_path)

slides_data = [data.get("slide1"), data.get("slide2"), data.get("slide3")]

for slide, img_data in zip(prs.slides, slides_data):
    if not img_data:
        continue

    img_bytes = base64.b64decode(img_data.split(",", 1)[1])
    img_stream = io.BytesIO(img_bytes)
    Image.open(img_stream)

    for shape in slide.shapes:
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height

            sp = shape._element
            sp.getparent().remove(sp)

            img_stream.seek(0)
            slide.shapes.add_picture(img_stream, left, top, width, height)
            break

prs.save(output_pptx)
print(f"PPTX gerado: {output_pptx}")
